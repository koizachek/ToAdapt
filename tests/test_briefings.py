"""Tests für die KI-Briefings (Master-Upload → Briefing je Stammgruppe).

Abgedeckt: Rubric-/Case-Config, Extraktion (PPTX mit Vorlagen-Shapes, DOCX,
PDF; Kenndaten; Zeichenzählung), formale Vorprüfung, Leitplanken-Nachprüfung,
Generator mit gemocktem LLM (valide Antwort, Garbage → technical_fallback,
Leitplanken-Treffer), Routen (fail-closed Auth, Master-Gate,
Sichtbarkeit je Übungsgruppe, DOCX-Download ohne Punkte, interne Einstufung
nur für Master, manuelle Zuordnung) und den Store-Datei-Fallback.

Alle Abgabetexte sind synthetisch — keine echten Teilnehmerdaten.
"""

import io
import json
import time
import zipfile
from datetime import date, timedelta

import pytest
from docx import Document
from fastapi.testclient import TestClient
from pptx import Presentation
from pptx.util import Inches

import backend.briefings.batches as batches_module
import backend.db.briefing_store as briefing_store_module
from backend.briefings import guardrails
from backend.briefings.batches import is_stale, new_batch
from backend.briefings.upload_token import UploadTokenError, sign_upload_token, verify_upload_token
from backend.briefings.extraction import (
    ZipValidationError,
    count_chars,
    extract_submission,
    iter_submission_entries,
    normalize_ueg,
    parse_code,
    parse_code_from_filename,
    split_bausteine,
)
from backend.briefings.formal import formal_checks, full_sentences_hint
from backend.briefings.generator import (
    BriefingGenerator,
    FALLBACK_TEXT,
    FeedbackGenerator,
    NO_CONTENT_TEXT,
    build_feedback_system_prompt,
    build_feedback_user_prompt,
    build_system_prompt,
    build_user_prompt,
)
from backend.briefings.rubrics import (
    SUPPORTED_TPS,
    case_context_for_tp,
    feedback_release_date,
    feedback_released,
    load_rubric,
    template_boilerplate,
)
from backend.llm import OpenRouterClient
from backend.main import app
from backend.timeutils import naive_utcnow

API_KEY = "tutor-key"

B1_TEXT = (
    "Die zwei kritischsten Herausforderungen sind der Kanalkonflikt und der auslaufende "
    "Patentschutz. Der Kanalkonflikt ist kritisch, weil der Direktvertrieb die Margen treibt, "
    "zugleich aber die Fachhändler bedroht. Wirkungskette: Der Patentschutz läuft aus, "
    "Wettbewerber bieten vergleichbare Systeme an, ON muss die Marke als zweiten Träger aufbauen."
)
B2_TEXT = (
    "Der Stakeholder, der den Handlungsspielraum am stärksten einschränkt, sind die Investoren. "
    "Ihre zentrale Erwartung ist profitables Wachstum. Die Implikation: ON kann die Distribution "
    "nicht verknappen, ohne einen glaubwürdigen Wachstumspfad zu zeigen."
)


# ---------------------------------------------------------------------------
# Synthetische Abgaben
# ---------------------------------------------------------------------------

def _template_pptx(tp: int, *, code: str, b1: str, b2: str, members: str = "Max Muster, Erika Beispiel") -> bytes:
    """Nachbau der offiziellen Vorlage: benannte KENN_*-Shapes auf Folie 1,
    KOPF_*-Vorlagentexte + Inhaltsplatzhalter auf Folie 2/3."""
    prs = Presentation()
    title_only = prs.slide_layouts[5]
    title_content = prs.slide_layouts[1]

    s1 = prs.slides.add_slide(title_only)
    s1.shapes.title.text = f"Touchpoint {tp} - Abgabe der Stammgruppe"
    for name, text in (
        ("L_UEG", "Übungsgruppe"), ("KENN_UEG", code.split("-")[1].replace("UEG", "")),
        ("L_SG", "Stammgruppe"), ("KENN_SG", code.split("-")[2].replace("SG", "")),
        ("L_TP", "Touchpoint"), ("KENN_TP", str(tp)),
        ("L_CODE", "Code"), ("KENN_CODE", code), ("H_CODE", "Format:  TP1-UEG07-SG3"),
        ("L_NAMEN", "Mitglieder der Stammgruppe"), ("KENN_NAMEN", members),
    ):
        box = s1.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(0.5))
        box.name = name
        box.text_frame.text = text

    for idx, (title, body) in enumerate((("Baustein 1 - Test", b1), ("Baustein 2 - Test", b2)), start=2):
        s = prs.slides.add_slide(title_content)
        s.shapes.title.text = title
        s.placeholders[1].text = body
        kopf = s.shapes.add_textbox(Inches(1), Inches(6), Inches(6), Inches(0.5))
        kopf.name = "KOPF_AUFTRAG"
        kopf.text_frame.text = "Beschreiben Sie zwei kritischsten Herausforderungen für ONs Geschäftsmodell"
        umfang = s.shapes.add_textbox(Inches(1), Inches(6.5), Inches(3), Inches(0.5))
        umfang.name = "KOPF_UMFANG"
        umfang.text_frame.text = "max. 1'350 Zeichen, min. 12 pt"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _docx(text_lines: list[str]) -> bytes:
    doc = Document()
    for line in text_lines:
        doc.add_paragraph(line)
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _minimal_pdf(text: str) -> bytes:
    stream = f"BT /F1 12 Tf 72 720 Td ({text}) Tj ET".encode("latin-1")
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
        b"/Contents 4 0 R /Resources << /Font << /F1 5 0 R >> >> >>",
        b"<< /Length " + str(len(stream)).encode() + b" >>\nstream\n" + stream + b"\nendstream",
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
    ]
    out = bytearray(b"%PDF-1.4\n")
    offsets = []
    for i, obj in enumerate(objects, start=1):
        offsets.append(len(out))
        out += f"{i} 0 obj\n".encode() + obj + b"\nendobj\n"
    xref = len(out)
    out += f"xref\n0 {len(objects) + 1}\n".encode() + b"0000000000 65535 f \n"
    for off in offsets:
        out += f"{off:010d} 00000 n \n".encode()
    out += f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\nstartxref\n{xref}\n%%EOF\n".encode()
    return bytes(out)


def _zip_of(files: dict[str, bytes]) -> bytes:
    buffer = io.BytesIO()
    with zipfile.ZipFile(buffer, "w") as archive:
        for name, data in files.items():
            archive.writestr(name, data)
    return buffer.getvalue()


def _llm_payload(**overrides) -> str:
    def baustein(prefix: str, names: list[str]):
        return {
            "kernposition": f"{prefix}: Die Gruppe hat sich für X entschieden, weil Y.",
            "tragende_argumente": [f"{prefix} Argument A", f"{prefix} Argument B", "Drittes wird gekappt"],
            "duenne_stellen": [f"{prefix}: Woran macht die Gruppe fest, dass …?"],
            "einschaetzung": f"{prefix}: Die Auswahl trägt, die Kette bleibt beim Mechanismus dünn.",
            "kriterien": [{"name": n, "niveau": "tragfaehig", "begruendung": "Weil."} for n in names],
        }
    rubric = load_rubric(1)
    data = {
        "baustein1": baustein("B1", [c.name for c in rubric.baustein("baustein1").criteria]),
        "baustein2": baustein("B2", [c.name for c in rubric.baustein("baustein2").criteria]),
        "judge_confidence": "high",
        "needs_human_review": False,
        "review_reason": None,
    }
    data.update(overrides)
    return json.dumps(data, ensure_ascii=False)


def _feedback_payload(**overrides) -> str:
    def baustein(prefix: str):
        return {
            "was_traegt": f"{prefix}: Ihre Auswahl ist am Fall belegt und die Erwartung konkret benannt.",
            "was_bleibt_duenn": f"{prefix}: Die Wirkungskette endet beim Umsatz; der Mechanismus davor fehlt.",
            "naechster_schritt": f"{prefix}: Formulieren Sie den Mechanismus zwischen Ursache und Umsatzfolge aus.",
        }
    data = {
        "baustein1": baustein("F1"),
        "baustein2": baustein("F2"),
        "feed_forward": "In Touchpoint 2 wird auf dieser Analyse entschieden; in der Klausur ist dies Aufgabe 1.",
        "judge_confidence": "high",
        "needs_human_review": False,
        "review_reason": None,
    }
    data.update(overrides)
    return json.dumps(data, ensure_ascii=False)


def _mock_llm(monkeypatch, response_text: str, feedback_text: str | None = None):
    """Antwortet auf den Briefing-Prompt mit response_text und auf den
    Feedback-Prompt (erkennbar am System-Prompt) mit feedback_text."""
    calls: list[dict] = []
    feedback_text = feedback_text if feedback_text is not None else _feedback_payload()

    async def fake_complete(self, *, system, messages, max_tokens, cache_system=False):
        calls.append({"system": system, "messages": messages, "cache_system": cache_system})
        if "Rückmeldung auf ihre Abgabe" in system:
            return feedback_text
        return response_text

    monkeypatch.setattr(OpenRouterClient, "complete", fake_complete)
    return calls


@pytest.fixture()
def client(monkeypatch, tmp_path):
    monkeypatch.setenv("TOADAPT_API_KEY", API_KEY)
    monkeypatch.setenv("OPENROUTER_API_KEY", "test-key")
    store_dir = tmp_path / "briefings"
    store_dir.mkdir()
    monkeypatch.setattr(briefing_store_module, "RESULTS_DIR", store_dir)
    batch_dir = tmp_path / "batches"
    batch_dir.mkdir()
    monkeypatch.setattr(batches_module, "BATCH_DIR", batch_dir)
    return TestClient(app)


def _master_headers() -> dict:
    return {"X-API-Key": API_KEY, "X-Teacher-Id": "master", "X-Teacher-Master": "1"}


def _tutor_headers(ueg: str) -> dict:
    return {"X-API-Key": API_KEY, "X-Teacher-Id": ueg, "X-Teacher-Master": "0"}


def _upload(client, files: dict[str, bytes], tp: int = 1, headers: dict | None = None, sync: bool = True):
    data = {"target_tp": str(tp)}
    if sync:
        data["sync"] = "1"
    return client.post(
        "/briefings/upload",
        files={"file": ("abgaben.zip", _zip_of(files), "application/zip")},
        data=data,
        headers=headers if headers is not None else _master_headers(),
    )


# ---------------------------------------------------------------------------
# Config
# ---------------------------------------------------------------------------

def test_all_rubrics_load_with_two_bausteine_and_examples():
    for tp in SUPPORTED_TPS:
        rubric = load_rubric(tp)
        assert rubric.tp == tp
        assert [b.key for b in rubric.bausteine] == ["baustein1", "baustein2"]
        assert all(b.criteria for b in rubric.bausteine)
        assert set(rubric.examples) == {"ueberzeugend", "tragfaehig", "ansatzweise"}
        assert "no_points_or_grades" in rubric.guardrails
        assert rubric.max_chars("baustein1") > 0 and rubric.max_chars("baustein2") > 0
        assert case_context_for_tp(tp)
        assert template_boilerplate(tp)


def test_tp2_case_context_includes_patent_section_from_chapter_a():
    ctx = case_context_for_tp(2)
    assert "# 3 Kapitel B" in ctx
    assert "## 2.8" in ctx


def test_unknown_tp_rejected():
    with pytest.raises(ValueError):
        load_rubric(6)


# ---------------------------------------------------------------------------
# Extraktion
# ---------------------------------------------------------------------------

def test_normalize_ueg_and_codes():
    assert normalize_ueg("7") == "UEG07"
    assert normalize_ueg("UEG07") == "UEG07"
    assert normalize_ueg("ueg 12") == "UEG12"
    assert normalize_ueg("Tutor A") == ""
    assert parse_code("Code: TP1-UEG07-SG3") == (1, "UEG07", 3)
    assert parse_code("tp2_ueg 3_sg8") == (2, "UEG03", 8)
    assert parse_code("kein code") is None
    assert parse_code_from_filename("TP1_UEG07_SG3.pptx") == (1, "UEG07", 3)


def test_pptx_template_extraction_reads_kenndaten_and_strips_boilerplate():
    data = _template_pptx(1, code="TP1-UEG07-SG3", b1=B1_TEXT, b2=B2_TEXT)
    sub = extract_submission("TP1_UEG07_SG3.pptx", data, 1)
    assert sub.format == "pptx" and sub.slide_count == 3 and sub.template_detected
    assert sub.kenndaten.code == "TP1-UEG07-SG3"
    assert sub.kenndaten.source == "kenndaten"
    assert sub.kenndaten.members_filled is True
    assert sub.baustein1 == B1_TEXT and sub.baustein2 == B2_TEXT
    assert "Beschreiben Sie" not in sub.baustein1
    assert "Zeichen" not in sub.baustein1
    assert sub.baustein1_chars == count_chars(B1_TEXT)
    # Mitgliedernamen dürfen nirgends im Ergebnis auftauchen
    assert "Max Muster" not in json.dumps(sub.__dict__, default=str)


def test_empty_template_is_not_assigned_to_example_code():
    """Die Vorlage trägt 'Format:  TP1-UEG07-SG3' als Beispiel — eine
    unausgefüllte Abgabe darf daraus keine Zuordnung ableiten."""
    prs = Presentation()
    s1 = prs.slides.add_slide(prs.slide_layouts[5])
    s1.shapes.title.text = "Touchpoint 1 - Abgabe der Stammgruppe"
    for name, text in (
        ("KENN_CODE", "[ bitte ausfüllen ]"), ("H_CODE", "Format:  TP1-UEG07-SG3"),
        ("DECK_FUSS", "Abgabe als .pptx … Dateiname: TP1_UEG07_SG3.pptx (analog zum Code)"),
        ("KENN_NAMEN", "[ bitte ausfüllen ]"),
    ):
        box = s1.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(0.5))
        box.name = name
        box.text_frame.text = text
    for _ in range(2):
        prs.slides.add_slide(prs.slide_layouts[1])
    buf = io.BytesIO()
    prs.save(buf)
    sub = extract_submission("abgabe.pptx", buf.getvalue(), 1)
    assert sub.kenndaten.code == "" and sub.kenndaten.source == ""
    assert sub.kenndaten.members_filled is False
    assert not sub.has_content
    # Fliesstext-Variante (DOCX/PDF): Format-Hinweis ebenfalls ignoriert
    assert parse_code("Format: TP1-UEG07-SG3\nDateiname: TP1_UEG07_SG3.pptx") is None
    assert parse_code("Format: TP1-UEG07-SG3\nCode: TP1-UEG12-SG4") == (1, "UEG12", 4)


def test_pptx_falls_back_to_filename_when_kenndaten_missing():
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[5])
    s.shapes.title.text = "Irgendein Deckblatt"
    for body in (B1_TEXT, B2_TEXT):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.placeholders[1].text = body
    buf = io.BytesIO()
    prs.save(buf)
    sub = extract_submission("TP1_UEG12_SG5.pptx", buf.getvalue(), 1)
    assert sub.kenndaten.source == "filename"
    assert sub.kenndaten.ueg == "UEG12" and sub.kenndaten.sg == 5
    assert not sub.template_detected


def test_docx_split_by_markers():
    data = _docx(["Code: TP1-UEG02-SG1", "Baustein 1 – Herausforderungen", B1_TEXT, "Baustein 2", B2_TEXT])
    sub = extract_submission("abgabe.docx", data, 1)
    assert sub.format == "docx"
    assert sub.kenndaten.code == "TP1-UEG02-SG1" and sub.kenndaten.source == "text"
    assert sub.baustein1 == B1_TEXT and sub.baustein2 == B2_TEXT


def test_flat_text_without_markers_lands_in_baustein1_with_note():
    sub = extract_submission("TP1_UEG02_SG2.pdf", _minimal_pdf("Alles in einem Absatz ohne Marker."), 1)
    assert sub.baustein1.startswith("Alles") and sub.baustein2 == ""
    assert any("Marker" in n for n in sub.notes)


def test_split_bausteine_variants():
    assert split_bausteine("Baustein 1\nA\nFolie 3 - Baustein 2\nB") == ("A", "B", True)
    assert split_bausteine("nur text") == ("nur text", "", False)


def test_zip_validation():
    with pytest.raises(ZipValidationError):
        list(iter_submission_entries(b"kein zip"))
    with pytest.raises(ZipValidationError):
        list(iter_submission_entries(_zip_of({"notes.txt": b"x"})))
    entries = list(iter_submission_entries(_zip_of({
        "__MACOSX/._a.pptx": b"junk", "ordner/TP1_UEG01_SG1.docx": _docx(["x"]), ".DS_Store": b"",
    })))
    assert [name for name, _ in entries] == ["TP1_UEG01_SG1.docx"]


def test_unreadable_file_raises():
    with pytest.raises(ValueError):
        extract_submission("kaputt.pptx", b"nicht pptx", 1)


# ---------------------------------------------------------------------------
# Formale Vorprüfung
# ---------------------------------------------------------------------------

def test_formal_checks_report_limits_and_patterns():
    rubric = load_rubric(1)
    long_text = "Satz. " * 300  # > 1'350 Zeichen
    sub = extract_submission("TP1_UEG07_SG3.pptx", _template_pptx(1, code="TP1-UEG07-SG3", b1=long_text, b2=B2_TEXT), 1)
    formal = formal_checks(sub, rubric, 1)
    assert formal["baustein1_within_limit"] is False and formal["baustein2_within_limit"] is True
    assert formal["code_valid"] is True and formal["filename_valid"] is True
    assert formal["code_matches_tp"] is True

    sub_wrong = extract_submission("abgabe.docx", _docx(["TP2-UEG07-SG3", "Baustein 1", "- Stichpunkt", "- noch einer", "Baustein 2", B2_TEXT]), 1)
    formal_wrong = formal_checks(sub_wrong, rubric, 1)
    assert formal_wrong["code_matches_tp"] is False
    assert formal_wrong["filename_valid"] is False
    assert any("Touchpoint 2" in n for n in formal_wrong["notes"])
    assert formal_wrong["full_sentences_hint"]


def test_full_sentences_hint():
    assert full_sentences_hint("Ein Satz.\nNoch ein Satz.") is None
    assert full_sentences_hint("Stichpunkt\nnoch einer\ndritter") is not None


# ---------------------------------------------------------------------------
# Leitplanken
# ---------------------------------------------------------------------------

@pytest.mark.parametrize("text,label", [
    ("Das ergibt 12 Punkte von 20.", "points"),
    ("Dafür erhält die Gruppe volle Punkte.", "points"),
    ("Die Note wäre gut.", "grades"),
    ("Niveau: tragfähig bei der Auswahl.", "scale"),
    ("Die richtige Entscheidung wäre der Fachhandel gewesen.", "model_solution"),
    ("Die Gruppe hätte die Investoren wählen sollen.", "model_solution"),
    ("Im Vergleich zu den anderen Gruppen ist das dünn.", "group_comparison"),
])
def test_guardrail_hits(text, label):
    assert label in guardrails.check_briefing_text(text)


@pytest.mark.parametrize("text", [
    "An diesem Punkt bleibt die Begründung dünn.",
    "Die Begründung trägt bei der Auswahl, bleibt aber bei der Wirkungskette dünn.",
    "Der Fachhandel erreicht 60 Prozent der Kunden.",
    "Woran macht die Gruppe fest, dass der Patentablauf kritischer ist als der Kanalkonflikt?",
    "Die Gruppe zählt vier Punkte aus Abschnitt 2.5 auf, ohne zu gewichten.",
])
def test_guardrail_allows_prose(text):
    assert guardrails.check_briefing_text(text) == []


def test_apply_guardrails_replaces_and_swissifies():
    cleaned, hits = guardrails.apply_guardrails("Die Gruppe erhält 5 Punkte.")
    assert cleaned == guardrails.GUARDRAIL_PLACEHOLDER and hits == ["points"]
    cleaned, hits = guardrails.apply_guardrails(["Grosse Straße.", "Musterlösung wäre X."])
    assert cleaned[0] == "Grosse Strasse." and cleaned[1] == guardrails.GUARDRAIL_PLACEHOLDER
    assert hits == ["model_solution"]


# ---------------------------------------------------------------------------
# Generator
# ---------------------------------------------------------------------------

def _sub(tp=1, b1=B1_TEXT, b2=B2_TEXT):
    return extract_submission("TP1_UEG07_SG3.pptx", _template_pptx(tp, code=f"TP{tp}-UEG07-SG3", b1=b1, b2=b2), tp)


def test_system_prompt_is_cached_and_contains_rubric_case_examples(monkeypatch):
    rubric = load_rubric(1)
    system = build_system_prompt(rubric)
    assert build_system_prompt(rubric) == system  # byte-identisch → Prompt-Caching
    assert "Auswahl und Kritikalität" in system and "## 2.5" in system
    assert "Beispielabgabe · tragfaehig" in system
    assert "Keine Punkte" in system and "Keine Musterlösung" in system
    user = build_user_prompt(rubric, _sub())
    assert B1_TEXT in user and "TP1-UEG07-SG3" in user


async def test_generator_valid_payload_splits_briefing_and_assessment(monkeypatch):
    calls = _mock_llm(monkeypatch, _llm_payload())
    result = await BriefingGenerator("k").generate(briefing_id="b1", rubric=load_rubric(1), sub=_sub())
    assert result["evaluation_status"] == "ok" and result["needs_human_review"] is False
    assert calls[0]["cache_system"] is True
    b1 = result["briefing"]["baustein1"]
    assert b1["kernposition"].startswith("B1")
    assert len(b1["tragende_argumente"]) == 2   # max 2 erzwungen
    assert "kriterien" not in b1                 # intern bleibt intern
    assert result["assessment"]["baustein1"]["kriterien"][0]["niveau"] == "tragfaehig"
    assert result["assessment"]["baustein1"]["fehlende_kriterien"] == []


async def test_generator_garbage_then_repair_then_fallback(monkeypatch):
    _mock_llm(monkeypatch, "das ist kein json")
    result = await BriefingGenerator("k").generate(briefing_id="b2", rubric=load_rubric(1), sub=_sub())
    assert result["evaluation_status"] == "technical_fallback"
    assert result["needs_human_review"] is True
    assert result["briefing"]["baustein1"]["kernposition"] == FALLBACK_TEXT


async def test_generator_transport_error_falls_back(monkeypatch):
    async def boom(self, **kwargs):
        raise RuntimeError("timeout")
    monkeypatch.setattr(OpenRouterClient, "complete", boom)
    result = await BriefingGenerator("k").generate(briefing_id="b3", rubric=load_rubric(1), sub=_sub())
    assert result["evaluation_status"] == "technical_fallback"


async def test_generator_guardrail_hit_replaces_field_and_flags(monkeypatch):
    payload = json.loads(_llm_payload())
    payload["baustein2"]["einschaetzung"] = "Die richtige Entscheidung wäre der Fachhandel gewesen."
    _mock_llm(monkeypatch, json.dumps(payload, ensure_ascii=False))
    result = await BriefingGenerator("k").generate(briefing_id="b4", rubric=load_rubric(1), sub=_sub())
    assert result["guardrail_hits"] == ["model_solution"]
    assert result["needs_human_review"] is True
    assert result["briefing"]["baustein2"]["einschaetzung"] == guardrails.GUARDRAIL_PLACEHOLDER
    assert result["briefing"]["baustein1"]["einschaetzung"].startswith("B1")


async def test_generator_empty_baustein_gets_placeholder_without_llm_for_no_content(monkeypatch):
    calls = _mock_llm(monkeypatch, _llm_payload())
    result = await BriefingGenerator("k").generate(briefing_id="b5", rubric=load_rubric(1), sub=_sub(b1=B1_TEXT, b2=""))
    assert result["briefing"]["baustein2"]["kernposition"] == NO_CONTENT_TEXT
    assert result["assessment"]["baustein2"]["keine_abgabe"] is True
    assert len(calls) == 1
    result = await BriefingGenerator("k").generate(briefing_id="b6", rubric=load_rubric(1), sub=_sub(b1="", b2=""))
    assert result["evaluation_status"] == "no_content" and len(calls) == 1


# ---------------------------------------------------------------------------
# Routen
# ---------------------------------------------------------------------------

def test_auth_fail_closed_and_wrong_key(client, monkeypatch):
    monkeypatch.delenv("TOADAPT_API_KEY")
    assert client.get("/briefings").status_code == 503
    monkeypatch.setenv("TOADAPT_API_KEY", API_KEY)
    assert client.get("/briefings", headers={"X-API-Key": "falsch"}).status_code == 401


def test_upload_requires_master(client, monkeypatch):
    _mock_llm(monkeypatch, _llm_payload())
    files = {"TP1_UEG07_SG3.pptx": _template_pptx(1, code="TP1-UEG07-SG3", b1=B1_TEXT, b2=B2_TEXT)}
    assert _upload(client, files, headers=_tutor_headers("UEG07")).status_code == 403
    # Direkt mit API-Key ohne Identitäts-Header = Operator (Skript) → erlaubt
    assert _upload(client, files, headers={"X-API-Key": API_KEY}).status_code == 202


def test_upload_flow_visibility_docx_and_assessment(client, monkeypatch):
    _mock_llm(monkeypatch, _llm_payload())
    files = {
        "TP1_UEG07_SG3.pptx": _template_pptx(1, code="TP1-UEG07-SG3", b1=B1_TEXT, b2=B2_TEXT),
        "TP1_UEG07_SG5.docx": _docx(["TP1-UEG07-SG5", "Baustein 1", B1_TEXT, "Baustein 2", B2_TEXT]),
        "TP1_UEG08_SG1.pptx": _template_pptx(1, code="TP1-UEG08-SG1", b1=B1_TEXT, b2=B2_TEXT),
        "ohne_code.docx": _docx(["Baustein 1", B1_TEXT, "Baustein 2", B2_TEXT]),
        "kaputt.pdf": b"kein pdf",
    }
    resp = _upload(client, files)
    assert resp.status_code == 202, resp.text
    body = resp.json()
    assert body["status"] == "done" and body["total"] == 5 and body["processed"] == 5
    assert body["briefed"] == 4 and body["failed"] == 1 and body["unassigned"] == 1
    assert all("assessment" not in b for b in body["briefings"])
    by_name = {b["filename"]: b for b in body["briefings"]}
    assert by_name["TP1_UEG07_SG3.pptx"]["code"] == "TP1-UEG07-SG3"
    assert by_name["TP1_UEG07_SG3.pptx"]["formal"]["baustein1_within_limit"] is True
    assert by_name["ohne_code.docx"]["needs_human_review"] is True
    assert by_name["kaputt.pdf"]["status"] == "extraction_failed"

    # ÜGL UEG07 sieht nur die eigenen zwei Briefings
    mine = client.get("/briefings", headers=_tutor_headers("UEG07")).json()
    assert sorted(b["code"] for b in mine) == ["TP1-UEG07-SG3", "TP1-UEG07-SG5"]
    assert client.get("/briefings", headers=_tutor_headers("UEG09")).json() == []
    assert client.get("/briefings", headers=_tutor_headers("Tutor Ohne Nummer")).json() == []
    # Master sieht alles (inkl. nicht zuordenbar + Extraktionsfehler)
    assert len(client.get("/briefings", headers=_master_headers()).json()) == 5

    # Übersicht: UEG07 hat 2 von 8, fehlende Stammgruppen gelistet
    overview = client.get("/briefings/overview?tp=1", headers=_tutor_headers("UEG07")).json()
    assert len(overview) == 1 and overview[0]["ueg"] == "UEG07"
    assert overview[0]["missing_groups"] == [1, 2, 4, 6, 7, 8]

    # Einzel-Briefing: fremde ÜGL bekommt 404, interne Einstufung nur Master
    foreign = by_name["TP1_UEG08_SG1.pptx"]["briefing_id"]
    assert client.get(f"/briefings/{foreign}", headers=_tutor_headers("UEG07")).status_code == 404
    assert client.get(f"/briefings/{foreign}/assessment", headers=_tutor_headers("UEG08")).status_code == 403
    assessment = client.get(f"/briefings/{foreign}/assessment", headers=_master_headers()).json()
    assert assessment["assessment"]["baustein1"]["kriterien"]

    # DOCX-Bundle für die eigene Übungsgruppe: echtes DOCX, ohne Punkte/Stufen
    docx_resp = client.get("/briefings/docx?tp=1", headers=_tutor_headers("UEG07"))
    assert docx_resp.status_code == 200
    assert docx_resp.headers["content-type"].startswith("application/vnd.openxmlformats")
    assert "KI-Briefing_TP1_UEG07.docx" in docx_resp.headers["content-disposition"]
    doc = Document(io.BytesIO(docx_resp.content))
    text = "\n".join(p.text for p in doc.paragraphs) + "\n".join(
        c.text for t in doc.tables for r in t.rows for c in r.cells
    )
    assert "Stammgruppe SG3" in text and "Stammgruppe SG5" in text
    assert "Keine Abgabe eingegangen: SG1, SG2, SG4, SG6, SG7, SG8" in text
    assert "B1: Die Gruppe hat sich" in text
    lowered = text.lower()
    assert "tragfaehig" not in lowered and "niveau" not in lowered and "punkte von" not in lowered
    assert "Weil." not in text  # interne Kriterien-Begründung bleibt intern
    # Master braucht ueg; fremde ÜGL bekommt für UEG08 nichts über tp-Filter hinaus
    assert client.get("/briefings/docx?tp=1", headers=_master_headers()).status_code == 422
    assert client.get("/briefings/docx?tp=1&ueg=8", headers=_master_headers()).status_code == 200
    single = client.get(f"/briefings/{foreign}/docx", headers=_master_headers())
    assert single.status_code == 200 and "TP1-UEG08-SG1" in single.headers["content-disposition"]

    # Manuelle Zuordnung des Datensatzes ohne Code
    unassigned = by_name["ohne_code.docx"]["briefing_id"]
    assert client.patch(f"/briefings/{unassigned}", json={"ueg": "7", "sg": 2}, headers=_tutor_headers("UEG07")).status_code == 403
    patched = client.patch(f"/briefings/{unassigned}", json={"ueg": "7", "sg": 2}, headers=_master_headers())
    assert patched.status_code == 200 and patched.json()["code"] == "TP1-UEG07-SG2"
    assert patched.json()["needs_human_review"] is False
    mine = client.get("/briefings?tp=1", headers=_tutor_headers("UEG07")).json()
    assert sorted(b["code"] for b in mine) == ["TP1-UEG07-SG2", "TP1-UEG07-SG3", "TP1-UEG07-SG5"]


def test_reupload_same_group_latest_wins(client, monkeypatch):
    _mock_llm(monkeypatch, _llm_payload())
    files = {"TP1_UEG07_SG3.pptx": _template_pptx(1, code="TP1-UEG07-SG3", b1=B1_TEXT, b2=B2_TEXT)}
    first = _upload(client, files).json()["briefings"][0]["briefing_id"]
    second = _upload(client, files).json()["briefings"][0]["briefing_id"]
    listed = client.get("/briefings?tp=1", headers=_tutor_headers("UEG07")).json()
    assert [b["briefing_id"] for b in listed] == [second]
    assert first != second


def test_upload_rejects_bad_tp_and_bad_zip(client, monkeypatch):
    _mock_llm(monkeypatch, _llm_payload())
    assert _upload(client, {"a.docx": _docx(["x"])}, tp=6).status_code == 422
    resp = client.post(
        "/briefings/upload",
        files={"file": ("x.zip", b"kein zip", "application/zip")},
        data={"target_tp": "1"},
        headers=_master_headers(),
    )
    assert resp.status_code == 400


def test_technical_fallback_is_flagged_in_record(client, monkeypatch):
    _mock_llm(monkeypatch, "{{{ garbage")
    files = {"TP1_UEG07_SG3.pptx": _template_pptx(1, code="TP1-UEG07-SG3", b1=B1_TEXT, b2=B2_TEXT)}
    body = _upload(client, files).json()
    assert body["briefings"][0]["evaluation_status"] == "technical_fallback"
    assert body["review"] == 1


def test_async_upload_returns_running_batch_and_finishes(client, monkeypatch):
    _mock_llm(monkeypatch, _llm_payload())
    files = {"TP1_UEG07_SG3.pptx": _template_pptx(1, code="TP1-UEG07-SG3", b1=B1_TEXT, b2=B2_TEXT)}
    # Hintergrund-Task braucht einen persistenten Event-Loop → TestClient als
    # Kontextmanager (ausserhalb des with-Blocks stirbt der Loop pro Request).
    with TestClient(app) as running:
        resp = _upload(running, files, sync=False)
        assert resp.status_code == 202
        body = resp.json()
        assert body["status"] in ("running", "done") and body["total"] == 1 and body["briefings"] == []
        batch_id = body["batch_id"]
        status = body
        for _ in range(100):
            status = running.get(f"/briefings/batches/{batch_id}", headers=_master_headers()).json()
            if status["status"] == "done":
                break
            time.sleep(0.05)
        assert status["status"] == "done" and status["processed"] == 1 and status["briefed"] == 1
        assert status["stale"] is False
        listed = running.get("/briefings/batches?tp=1", headers=_master_headers()).json()
        assert [b["batch_id"] for b in listed] == [batch_id]
        assert running.get("/briefings/batches", headers=_tutor_headers("UEG07")).status_code == 403
        assert running.get("/briefings?tp=1", headers=_tutor_headers("UEG07")).json()[0]["code"] == "TP1-UEG07-SG3"


def test_stale_batch_flag():
    batch = new_batch(batch_id="b", target_tp=1, total=3, uploaded_by=None, filename="x.zip")
    assert is_stale(batch) is False
    batch["updated_at"] = (naive_utcnow() - timedelta(hours=1)).isoformat()
    assert is_stale(batch) is True
    batch["status"] = "done"
    assert is_stale(batch) is False


def test_upload_token_auth(client, monkeypatch):
    _mock_llm(monkeypatch, _llm_payload())
    files = {"TP1_UEG07_SG3.pptx": _template_pptx(1, code="TP1-UEG07-SG3", b1=B1_TEXT, b2=B2_TEXT)}
    token = sign_upload_token(tutor="master", master=True)
    resp = _upload(client, files, headers={"X-Upload-Token": token})
    assert resp.status_code == 202 and resp.json()["uploaded_by"] == "master"
    # Nicht-Master-Token, manipuliertes Token, abgelaufenes Token, fehlendes Token
    assert _upload(client, files, headers={"X-Upload-Token": sign_upload_token(tutor="UEG07", master=False)}).status_code == 401
    assert _upload(client, files, headers={"X-Upload-Token": token[:-3] + "abc"}).status_code == 401
    expired = sign_upload_token(tutor="master", master=True, ttl_seconds=-120)
    assert _upload(client, files, headers={"X-Upload-Token": expired}).status_code == 401
    assert _upload(client, files, headers={"X-Nothing": "1"}).status_code == 401
    # Token nur auf der Upload-Route gültig — Lese-Routen verlangen den API-Key
    assert client.get("/briefings", headers={"X-Upload-Token": token}).status_code == 401
    # Fail-closed ohne konfigurierten Key
    monkeypatch.delenv("TOADAPT_API_KEY")
    assert _upload(client, files, headers={"X-Upload-Token": token}).status_code == 503


def test_verify_upload_token_roundtrip(monkeypatch):
    monkeypatch.setenv("TOADAPT_API_KEY", API_KEY)
    payload = verify_upload_token(sign_upload_token(tutor="master", master=True, jti="j-1"))
    assert payload["tutor"] == "master" and payload["jti"] == "j-1"
    with pytest.raises(UploadTokenError):
        verify_upload_token("kaputt")
    token = sign_upload_token(tutor="master", master=True)
    monkeypatch.setenv("TOADAPT_API_KEY", "anderer-key")
    with pytest.raises(UploadTokenError):
        verify_upload_token(token)


# ---------------------------------------------------------------------------
# Produkt 2: KI-Feedback (Freigabe erst nach dem Termin)
# ---------------------------------------------------------------------------

def test_feedback_release_date_is_day_after_termin(monkeypatch):
    assert feedback_release_date(1) == date(2026, 10, 3)
    assert feedback_released(1, today=date(2026, 10, 2)) is False
    assert feedback_released(1, today=date(2026, 10, 3)) is True
    assert feedback_release_date(9) is None and feedback_released(9) is False


async def test_feedback_generator_valid_and_guardrail(monkeypatch):
    _mock_llm(monkeypatch, _llm_payload())
    gen = FeedbackGenerator("k")
    result = await gen.generate_feedback(briefing_id="f1", rubric=load_rubric(1), sub=_sub(), assessment=None)
    assert result["feedback_status"] == "ok" and result["feedback_needs_human_review"] is False
    assert result["feedback"]["baustein1"]["naechster_schritt"].startswith("F1")
    assert result["feedback"]["feed_forward"].startswith("In Touchpoint 2")
    # Guardrail: Musterlösung im Feedback → Platzhalter + Review
    payload = json.loads(_feedback_payload())
    payload["baustein2"]["naechster_schritt"] = "Die richtige Entscheidung wäre der Fachhandel gewesen."
    _mock_llm(monkeypatch, _llm_payload(), json.dumps(payload, ensure_ascii=False))
    result = await gen.generate_feedback(briefing_id="f2", rubric=load_rubric(1), sub=_sub(), assessment=None)
    assert result["feedback_guardrail_hits"] == ["model_solution"]
    assert result["feedback"]["baustein2"]["naechster_schritt"] == guardrails.GUARDRAIL_PLACEHOLDER
    assert result["feedback_needs_human_review"] is True
    # Garbage → technical_fallback mit Feed-forward-Anker
    _mock_llm(monkeypatch, _llm_payload(), "kein json")
    result = await gen.generate_feedback(briefing_id="f3", rubric=load_rubric(1), sub=_sub(), assessment=None)
    assert result["feedback_status"] == "technical_fallback"
    assert result["feedback"]["baustein1"]["was_traegt"] == FALLBACK_TEXT
    assert "Touchpoint 2" in result["feedback"]["feed_forward"]


def test_feedback_prompt_contains_anchor_and_assessment():
    rubric = load_rubric(1)
    system = build_feedback_system_prompt(rubric)
    assert build_feedback_system_prompt(rubric) == system
    assert "Rückmeldung auf ihre Abgabe" in system and "Aufgabe 1" in system
    assessment = {"baustein1": {"kriterien": [{"name": "Erläuterung", "niveau": "tragfaehig", "begruendung": "x"}]}}
    user = build_feedback_user_prompt(rubric, _sub(), assessment)
    assert "Erläuterung: tragfaehig" in user


def test_feedback_gating_and_downloads(client, monkeypatch):
    import backend.briefings.routes as routes_module

    _mock_llm(monkeypatch, _llm_payload())
    files = {
        "TP1_UEG07_SG3.pptx": _template_pptx(1, code="TP1-UEG07-SG3", b1=B1_TEXT, b2=B2_TEXT),
        "TP1_UEG07_SG5.pptx": _template_pptx(1, code="TP1-UEG07-SG5", b1=B1_TEXT, b2=B2_TEXT),
    }
    # Vor dem Termin
    monkeypatch.setattr(routes_module, "feedback_released", lambda tp: False)
    body = _upload(client, files).json()
    rec = body["briefings"][0]
    assert rec["feedback_status"] == "ok" and rec["feedback_released"] is False
    assert rec["feedback_available_from"] == "2026-10-03"
    assert rec["feedback"]["baustein1"]["was_traegt"].startswith("F1")   # Master sieht Inhalt
    bid = rec["briefing_id"]
    mine = client.get("/briefings?tp=1", headers=_tutor_headers("UEG07")).json()
    assert all(b["feedback"] == {} for b in mine)                          # ÜGL nicht vor Freigabe
    assert client.get(f"/briefings/{bid}/feedback/docx", headers=_tutor_headers("UEG07")).status_code == 423
    assert client.get("/briefings/feedback/zip?tp=1", headers=_tutor_headers("UEG07")).status_code == 423
    assert client.get(f"/briefings/{bid}/feedback/docx?force=1", headers=_tutor_headers("UEG07")).status_code == 423
    forced = client.get(f"/briefings/{bid}/feedback/docx?force=1", headers=_master_headers())
    assert forced.status_code == 200                                        # Master-QS, geloggt

    # Nach dem Termin
    monkeypatch.setattr(routes_module, "feedback_released", lambda tp: True)
    mine = client.get("/briefings?tp=1", headers=_tutor_headers("UEG07")).json()
    assert all(b["feedback_released"] and b["feedback"]["feed_forward"] for b in mine)
    single = client.get(f"/briefings/{bid}/feedback/docx", headers=_tutor_headers("UEG07"))
    assert single.status_code == 200 and "KI-Feedback_TP1-UEG07-SG3.docx" in single.headers["content-disposition"]
    doc = Document(io.BytesIO(single.content))
    text = "\n".join(p.text for p in doc.paragraphs)
    assert "Stammgruppe SG3" in text and "Was trägt:" in text and "Nächster Schritt:" in text
    assert "Ausblick" in text and "In Touchpoint 2" in text
    lowered = text.lower()
    assert "tragfaehig" not in lowered and "niveau" not in lowered and "punkte von" not in lowered
    assert "Formale Vorprüfung" not in text and "Kernposition" not in text   # kein Briefing-Inhalt

    bundle = client.get("/briefings/feedback/zip?tp=1", headers=_tutor_headers("UEG07"))
    assert bundle.status_code == 200 and bundle.headers["content-type"] == "application/zip"
    names = sorted(zipfile.ZipFile(io.BytesIO(bundle.content)).namelist())
    assert names == ["KI-Feedback_TP1-UEG07-SG3.docx", "KI-Feedback_TP1-UEG07-SG5.docx"]
    assert client.get("/briefings/feedback/zip?tp=1", headers=_tutor_headers("UEG08")).status_code == 404
    assert client.get("/briefings/feedback/zip?tp=1", headers=_master_headers()).status_code == 422
    assert client.get("/briefings/feedback/zip?tp=1&ueg=7", headers=_master_headers()).status_code == 200


def test_store_file_fallback_roundtrip(monkeypatch, tmp_path):
    monkeypatch.setattr(briefing_store_module, "RESULTS_DIR", tmp_path)
    store = briefing_store_module.BriefingStore()
    store.save({"briefing_id": "x1", "target_tp": 1, "ueg": "UEG01", "sg": 1})
    assert (tmp_path / "x1.json").exists()
    assert store.get("x1")["ueg"] == "UEG01"
    assert store.get("nope") is None
