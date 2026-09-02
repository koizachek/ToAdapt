"""Extraktion der Stammgruppen-Abgaben (PPTX, DOCX, PDF) aus dem Master-ZIP.

Alles läuft rein im Speicher (nichts wird entpackt auf Platte geschrieben),
ohne LLM und ist pur testbar. Ergebnis je Datei ist ein
``ExtractedSubmission`` mit:

- Kenndaten (Touchpoint, Übungsgruppe UEGxx, Stammgruppe SGy, Code) — aus
  den benannten Kenndaten-Shapes der offiziellen Vorlage (``KENN_*``),
  ersatzweise per Regex aus dem Text, ersatzweise aus dem Dateinamen.
- Text je Baustein (Folie 2 = Baustein 1, Folie 3 = Baustein 2), bereinigt
  um Vorlagentext (Titel, Auftragstext, Umfangshinweis, Platzhalter).
- Zeichenzahl je Baustein (inklusive Leerzeichen, ohne Absatzmarken —
  wie PowerPoint/Word zählen).

Die Mitgliedernamen auf dem Deckblatt (``KENN_NAMEN``) werden NIE
übernommen — nur ob das Feld ausgefüllt ist (formale Vorprüfung).
"""

from __future__ import annotations

import io
import re
import zipfile
from dataclasses import dataclass, field
from typing import Iterator

from pypdf import PdfReader

from backend.briefings.rubrics import SUPPORTED_TPS, template_boilerplate

SUPPORTED_EXTENSIONS = (".pptx", ".docx", ".pdf")

# Limits gegen ZIP-Bomben und Speicherfrass. Ein Semester-Upload umfasst bis
# zu 55 Übungsgruppen × 8 Stammgruppen = 440 Dateien je Touchpoint; die
# Vorlage wiegt ~0,7 MB (Bilder), daher grosszügige Einzel-/Gesamtlimits.
MAX_ZIP_ENTRIES = 600
MAX_FILE_BYTES = 40 * 1024 * 1024
MAX_TOTAL_BYTES = 1200 * 1024 * 1024

# Deckel für den Judge-Input je Baustein (die Vorlage erlaubt max. 1'650
# Zeichen; deutlich mehr deutet auf Fehlextraktion oder Regelverstoss hin).
MAX_BAUSTEIN_CHARS = 8000

CODE_RE = re.compile(
    r"TP\s*([1-5])\s*[-_– ]\s*UEG\s*(\d{1,2})\s*[-_– ]\s*SG\s*([1-8])\b",
    re.IGNORECASE,
)
FILENAME_RE = re.compile(
    r"TP\s*([1-5])[-_ ]UEG\s*(\d{1,2})[-_ ]SG\s*([1-8])",
    re.IGNORECASE,
)
_BAUSTEIN_MARK = {
    1: re.compile(r"^\s*(?:folie\s*2\s*[-–:·]?\s*)?baustein\s*1\b", re.IGNORECASE),
    2: re.compile(r"^\s*(?:folie\s*3\s*[-–:·]?\s*)?baustein\s*2\b", re.IGNORECASE),
}
_PLACEHOLDER_PREFIXES = ("[ihre darstellung", "[ bitte ausfüllen", "[bitte ausfüllen")
_EMPTY_FIELD = {"", "__", "_", "[ bitte ausfüllen ]", "[bitte ausfüllen]"}


class ZipValidationError(ValueError):
    """Ungültiges oder zu grosses ZIP — als 400 an den Aufrufer."""


@dataclass
class Kenndaten:
    tp: int | None = None
    ueg: str = ""            # "UEG07"
    sg: int | None = None    # 1–8
    code: str = ""           # "TP1-UEG07-SG3"
    source: str = ""         # "kenndaten" | "text" | "filename" | ""
    members_filled: bool | None = None


@dataclass
class ExtractedSubmission:
    filename: str
    format: str                       # "pptx" | "docx" | "pdf"
    kenndaten: Kenndaten
    baustein1: str = ""
    baustein2: str = ""
    baustein1_chars: int = 0
    baustein2_chars: int = 0
    slide_count: int | None = None
    template_detected: bool = False
    notes: list[str] = field(default_factory=list)

    @property
    def has_content(self) -> bool:
        return bool(self.baustein1.strip() or self.baustein2.strip())


# ---------------------------------------------------------------------------
# Hilfsfunktionen
# ---------------------------------------------------------------------------

def normalize_ueg(value: str | int | None) -> str:
    """'7' / '07' / 'UEG07' / 'ueg 7' → 'UEG07'; '' wenn nicht erkennbar."""
    if value is None:
        return ""
    match = re.fullmatch(r"\s*(?:UEG)?\s*0*(\d{1,2})\s*", str(value), re.IGNORECASE)
    if not match:
        return ""
    return f"UEG{int(match.group(1)):02d}"


def parse_uegs(value: str | None) -> list[str]:
    """Alle Übungsgruppen einer Tutor-Kennung: 'UEG07' → ['UEG07'];
    'UEG07+UEG12' / 'UEG07, UEG12' / '7 12' → ['UEG07', 'UEG12'].
    Eine ÜGL kann mehrere Übungsgruppen führen — die Kennung (Schlüssel in
    TEACHER_ACCESS_CODES) trägt sie alle. Nicht parsebare Teile werden ignoriert."""
    out: list[str] = []
    for token in re.split(r"[^A-Za-z0-9]+", str(value or "")):
        ueg = normalize_ueg(token)
        if ueg and ueg not in out:
            out.append(ueg)
    return out


def build_code(tp: int, ueg: str, sg: int) -> str:
    return f"TP{tp}-{ueg}-SG{sg}"


def _norm(text: str) -> str:
    return re.sub(r"\s+", " ", text.replace("\x0b", " ")).strip().lower()


def count_chars(text: str) -> int:
    """Zeichen inklusive Leerzeichen, ohne Zeilen-/Absatzumbrüche."""
    return len(text.replace("\r", "").replace("\n", "").replace("\x0b", ""))


_FORMAT_HINT_RE = re.compile(
    r"(format|dateiname|beispiel|z\.?\s?b\.?|analog zum code)\s*[:：]?\s*"
    r"TP\s*[1-5]\s*[-_– ]\s*UEG\s*\d{1,2}\s*[-_– ]\s*SG\s*[1-8](\.\w+)?",
    re.IGNORECASE,
)


def _without_format_hints(text: str) -> str:
    """Entfernt Beispiel-Codes aus Vorlagentexten ("Format:  TP1-UEG07-SG3",
    "Dateiname: TP1_UEG07_SG3.pptx"), damit eine unausgefüllte Vorlage nicht
    der Beispielgruppe zugeordnet wird."""
    return _FORMAT_HINT_RE.sub(" ", text or "")


def parse_code(text: str) -> tuple[int, str, int] | None:
    match = CODE_RE.search(_without_format_hints(text))
    if not match:
        return None
    return int(match.group(1)), normalize_ueg(match.group(2)), int(match.group(3))


def parse_code_from_filename(filename: str) -> tuple[int, str, int] | None:
    match = FILENAME_RE.search(filename or "")
    if not match:
        return None
    return int(match.group(1)), normalize_ueg(match.group(2)), int(match.group(3))


def split_bausteine(text: str) -> tuple[str, str, bool]:
    """Teilt Fliesstext an 'Baustein 1' / 'Baustein 2'-Markern.

    Rückgabe (baustein1, baustein2, gefunden). Ohne Marker landet alles in
    Baustein 1 und ``gefunden`` ist False.
    """
    lines = (text or "").splitlines()
    idx = {1: None, 2: None}
    for i, line in enumerate(lines):
        for key, pattern in _BAUSTEIN_MARK.items():
            if idx[key] is None and pattern.match(line):
                idx[key] = i
    if idx[1] is None and idx[2] is None:
        return text.strip(), "", False
    start1 = (idx[1] + 1) if idx[1] is not None else 0
    if idx[2] is None:
        return "\n".join(lines[start1:]).strip(), "", True
    if idx[1] is None or idx[2] < idx[1]:
        return "", "\n".join(lines[idx[2] + 1:]).strip(), True
    b1 = "\n".join(lines[start1: idx[2]]).strip()
    b2 = "\n".join(lines[idx[2] + 1:]).strip()
    return b1, b2, True


def _strip_boilerplate_lines(text: str, boilerplate: set[str]) -> str:
    kept = []
    for line in (text or "").splitlines():
        n = _norm(line)
        if not n:
            continue
        if n in boilerplate or n.startswith(_PLACEHOLDER_PREFIXES):
            continue
        if re.fullmatch(r"\d{1,2}", n):          # Foliennummer
            continue
        if re.match(r"^max\.?\s*\d", n):          # Umfangshinweis
            continue
        kept.append(line.rstrip())
    return "\n".join(kept).strip()


def _boilerplate_set(tp: int | None) -> set[str]:
    tps = [tp] if tp in SUPPORTED_TPS else list(SUPPORTED_TPS)
    out: set[str] = set()
    for t in tps:
        out.update(_norm(s) for s in template_boilerplate(t))
    return out


# ---------------------------------------------------------------------------
# ZIP
# ---------------------------------------------------------------------------

def iter_submission_entries(zip_bytes: bytes) -> Iterator[tuple[str, bytes]]:
    """Liefert (dateiname, bytes) für alle PPTX/DOCX/PDF-Einträge, lazy.

    Verzeichnisse, macOS-Metadaten und andere Formate werden übersprungen.
    Wirft ZipValidationError bei kaputtem Archiv oder verletzten Limits.
    """
    try:
        archive = zipfile.ZipFile(io.BytesIO(zip_bytes))
    except zipfile.BadZipFile as exc:
        raise ZipValidationError("Datei ist kein gültiges ZIP-Archiv.") from exc

    infos = [
        info for info in archive.infolist()
        if not info.is_dir()
        and not info.filename.startswith("__MACOSX/")
        and not info.filename.rsplit("/", 1)[-1].startswith(".")
    ]
    if len(infos) > MAX_ZIP_ENTRIES:
        raise ZipValidationError(f"ZIP enthält zu viele Dateien (max. {MAX_ZIP_ENTRIES}).")

    selected = [i for i in infos if i.filename.lower().endswith(SUPPORTED_EXTENSIONS)]
    if not selected:
        raise ZipValidationError("ZIP enthält keine PPTX-, DOCX- oder PDF-Dateien.")

    total = 0
    for info in selected:
        if info.file_size > MAX_FILE_BYTES:
            raise ZipValidationError(
                f"'{info.filename}' überschreitet das Dateilimit von "
                f"{MAX_FILE_BYTES // (1024 * 1024)} MB."
            )
        total += info.file_size
        if total > MAX_TOTAL_BYTES:
            raise ZipValidationError("ZIP-Inhalt überschreitet das Gesamt-Limit.")

    for info in selected:
        yield info.filename.rsplit("/", 1)[-1], archive.read(info)


# ---------------------------------------------------------------------------
# PPTX
# ---------------------------------------------------------------------------

def _shape_texts(shape) -> list[tuple[str, str]]:
    """(shape_name, text) rekursiv über Gruppen und Tabellen."""
    out: list[tuple[str, str]] = []
    if getattr(shape, "shape_type", None) == 6 and hasattr(shape, "shapes"):  # GROUP
        for sub in shape.shapes:
            out.extend(_shape_texts(sub))
        return out
    if getattr(shape, "has_table", False) and shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                if cell.text.strip():
                    out.append((shape.name, cell.text))
        return out
    if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
        text = shape.text_frame.text
        if text.strip():
            out.append((shape.name, text))
    return out


def _is_pptx_boilerplate(shape, text: str, boilerplate: set[str]) -> bool:
    name = (getattr(shape, "name", "") or "").upper()
    if name.startswith("KOPF_") or name.startswith("L_") or name.startswith("H_"):
        return True
    if getattr(shape, "is_placeholder", False):
        ph_type = shape.placeholder_format.type
        if ph_type in (1, 3, 13, 15, 16):   # TITLE, CENTER_TITLE, SLIDE_NUMBER, FOOTER, HEADER
            return True
    n = _norm(text)
    return n in boilerplate or n.startswith(_PLACEHOLDER_PREFIXES) or bool(re.fullmatch(r"\d{1,2}", n))


def _slide_content(slide, boilerplate: set[str]) -> str:
    parts: list[str] = []
    for shape in slide.shapes:
        for _, text in _shape_texts(shape):
            if _is_pptx_boilerplate(shape, text, boilerplate):
                continue
            cleaned = _strip_boilerplate_lines(text, boilerplate)
            if cleaned:
                parts.append(cleaned)
    return "\n".join(parts).strip()


def _kenndaten_from_pptx(slide, filename: str) -> tuple[Kenndaten, bool]:
    """Liest die KENN_*-Shapes der Vorlage; Rückgabe (kenndaten, vorlage_erkannt)."""
    fields: dict[str, str] = {}
    for shape in slide.shapes:
        name = (getattr(shape, "name", "") or "").upper()
        if name.startswith("KENN_") and getattr(shape, "has_text_frame", False):
            fields[name] = shape.text_frame.text.strip()
    template = bool(fields)
    kd = Kenndaten()
    if "KENN_NAMEN" in fields:
        kd.members_filled = _norm(fields["KENN_NAMEN"]) not in {_norm(v) for v in _EMPTY_FIELD}

    parsed = parse_code(fields.get("KENN_CODE", ""))
    if parsed:
        kd.tp, kd.ueg, kd.sg = parsed
        kd.source = "kenndaten"
    else:
        ueg = normalize_ueg(fields.get("KENN_UEG", ""))
        sg_match = re.fullmatch(r"\s*0?([1-8])\s*", fields.get("KENN_SG", "") or "")
        tp_match = re.fullmatch(r"\s*([1-5])\s*", fields.get("KENN_TP", "") or "")
        if ueg and sg_match:
            kd.ueg, kd.sg = ueg, int(sg_match.group(1))
            kd.tp = int(tp_match.group(1)) if tp_match else None
            kd.source = "kenndaten"
    if not kd.source:
        # Nur Nicht-Vorlagen-Shapes: Labels (L_*), Hinweise (H_*) und die
        # Fusszeile (DECK_FUSS) tragen Beispielcodes, keine Kenndaten.
        all_text = "\n".join(
            t for shape in slide.shapes for name, t in _shape_texts(shape)
            if not name.upper().startswith(("L_", "H_", "DECK_", "KOPF_"))
        )
        parsed = parse_code(all_text)
        if parsed:
            kd.tp, kd.ueg, kd.sg = parsed
            kd.source = "text"
    if not kd.source:
        parsed = parse_code_from_filename(filename)
        if parsed:
            kd.tp, kd.ueg, kd.sg = parsed
            kd.source = "filename"
    if kd.ueg and kd.sg and kd.tp:
        kd.code = build_code(kd.tp, kd.ueg, kd.sg)
    return kd, template


def extract_pptx(filename: str, data: bytes, expected_tp: int | None = None) -> ExtractedSubmission:
    from pptx import Presentation  # lazy: schwerer Import

    try:
        prs = Presentation(io.BytesIO(data))
        slides = list(prs.slides)
    except Exception as exc:
        raise ValueError("PPTX konnte nicht gelesen werden.") from exc
    if not slides:
        raise ValueError("PPTX enthält keine Folien.")

    boilerplate = _boilerplate_set(expected_tp)
    kd, template = _kenndaten_from_pptx(slides[0], filename)
    sub = ExtractedSubmission(
        filename=filename, format="pptx", kenndaten=kd,
        slide_count=len(slides), template_detected=template,
    )

    if len(slides) >= 3:
        sub.baustein1 = _slide_content(slides[1], boilerplate)
        sub.baustein2 = _slide_content(slides[2], boilerplate)
        if len(slides) > 3:
            sub.notes.append(f"Abgabe hat {len(slides)} Folien (Vorlage: 3); nur Folie 2 und 3 wurden gelesen.")
    else:
        text = "\n".join(_slide_content(s, boilerplate) for s in slides)
        b1, b2, found = split_bausteine(text)
        sub.baustein1, sub.baustein2 = b1, b2
        sub.notes.append(
            f"Abgabe hat nur {len(slides)} Folie(n) (Vorlage: 3); Bausteine "
            + ("über Marker getrennt." if found else "konnten nicht getrennt werden.")
        )
    _finalize(sub)
    return sub


# ---------------------------------------------------------------------------
# DOCX
# ---------------------------------------------------------------------------

def _docx_text(data: bytes) -> str:
    from docx import Document  # lazy
    from docx.oxml.ns import qn

    try:
        doc = Document(io.BytesIO(data))
    except Exception as exc:
        raise ValueError("DOCX konnte nicht gelesen werden.") from exc

    lines: list[str] = []
    body = doc.element.body
    for child in body.iterchildren():
        if child.tag == qn("w:p"):
            text = "".join(t.text or "" for t in child.iter(qn("w:t")))
            if text.strip():
                lines.append(text)
        elif child.tag == qn("w:tbl"):
            for row in child.iter(qn("w:tr")):
                cells = []
                for cell in row.iter(qn("w:tc")):
                    ctext = " ".join(
                        "".join(t.text or "" for t in p.iter(qn("w:t")))
                        for p in cell.iter(qn("w:p"))
                    ).strip()
                    if ctext:
                        cells.append(ctext)
                if cells:
                    lines.append(" | ".join(cells))
    return "\n".join(lines)


def extract_docx(filename: str, data: bytes, expected_tp: int | None = None) -> ExtractedSubmission:
    text = _docx_text(data)
    if not text.strip():
        raise ValueError("DOCX enthält keinen Text.")
    return _from_flat_text(filename, "docx", text, expected_tp)


# ---------------------------------------------------------------------------
# PDF
# ---------------------------------------------------------------------------

def extract_pdf(filename: str, data: bytes, expected_tp: int | None = None) -> ExtractedSubmission:
    try:
        reader = PdfReader(io.BytesIO(data))
        pages = [(page.extract_text() or "").strip() for page in reader.pages]
    except Exception as exc:
        raise ValueError("PDF konnte nicht gelesen werden.") from exc
    if not any(pages):
        raise ValueError("PDF enthält keinen extrahierbaren Text (Scan ohne OCR?).")

    boilerplate = _boilerplate_set(expected_tp)
    if len(pages) >= 3:
        # Aus der Vorlage exportiert: Seite 2 = Baustein 1, Seite 3 = Baustein 2.
        kd = _kenndaten_from_text(pages[0], filename)
        sub = ExtractedSubmission(filename=filename, format="pdf", kenndaten=kd, slide_count=len(pages))
        sub.baustein1 = _strip_boilerplate_lines(pages[1], boilerplate)
        sub.baustein2 = _strip_boilerplate_lines(pages[2], boilerplate)
        if len(pages) > 3:
            sub.notes.append(f"PDF hat {len(pages)} Seiten (Vorlage: 3); nur Seite 2 und 3 wurden gelesen.")
        _finalize(sub)
        return sub
    return _from_flat_text(filename, "pdf", "\n".join(pages), expected_tp, page_count=len(pages))


# ---------------------------------------------------------------------------
# Gemeinsam
# ---------------------------------------------------------------------------

def _kenndaten_from_text(text: str, filename: str) -> Kenndaten:
    kd = Kenndaten()
    parsed = parse_code(text)
    if parsed:
        kd.tp, kd.ueg, kd.sg = parsed
        kd.source = "text"
    else:
        parsed = parse_code_from_filename(filename)
        if parsed:
            kd.tp, kd.ueg, kd.sg = parsed
            kd.source = "filename"
    if kd.ueg and kd.sg and kd.tp:
        kd.code = build_code(kd.tp, kd.ueg, kd.sg)
    return kd


def _from_flat_text(
    filename: str, fmt: str, text: str, expected_tp: int | None, page_count: int | None = None
) -> ExtractedSubmission:
    boilerplate = _boilerplate_set(expected_tp)
    kd = _kenndaten_from_text(text[:3000], filename)
    cleaned = _strip_boilerplate_lines(text, boilerplate)
    b1, b2, found = split_bausteine(cleaned)
    sub = ExtractedSubmission(filename=filename, format=fmt, kenndaten=kd, slide_count=page_count)
    sub.baustein1, sub.baustein2 = b1, b2
    if not found:
        sub.notes.append(
            "Keine 'Baustein 1'/'Baustein 2'-Marker gefunden; der gesamte Text wurde Baustein 1 zugeordnet."
        )
    _finalize(sub)
    return sub


def _finalize(sub: ExtractedSubmission) -> None:
    for key in ("baustein1", "baustein2"):
        text = getattr(sub, key)
        if len(text) > MAX_BAUSTEIN_CHARS:
            sub.notes.append(f"{key}: Text auf {MAX_BAUSTEIN_CHARS} Zeichen gekürzt.")
            text = text[:MAX_BAUSTEIN_CHARS]
            setattr(sub, key, text)
        setattr(sub, f"{key}_chars", count_chars(text))


def extract_submission(filename: str, data: bytes, expected_tp: int | None = None) -> ExtractedSubmission:
    """Dispatch nach Dateiendung; ValueError bei unlesbarer Datei."""
    lower = filename.lower()
    if lower.endswith(".pptx"):
        return extract_pptx(filename, data, expected_tp)
    if lower.endswith(".docx"):
        return extract_docx(filename, data, expected_tp)
    if lower.endswith(".pdf"):
        return extract_pdf(filename, data, expected_tp)
    raise ValueError("Nicht unterstütztes Format (erlaubt: pptx, docx, pdf).")
